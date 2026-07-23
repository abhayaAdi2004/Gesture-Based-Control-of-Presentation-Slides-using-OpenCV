"""
create_sample_ppt.py
--------------------
Creates a sample PowerPoint presentation for testing the gesture control system.

Usage:
    python create_sample_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1 - Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Title text
txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Gesture Controlled PowerPoint"
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Control your presentation with hand gestures!"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p2.alignment = PP_ALIGN.CENTER

# Slide 2 - Available Gestures
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.background
fill2 = bg2.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Hand Gesture Controls"
p3.font.size = Pt(36)
p3.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
p3.font.bold = True
p3.alignment = PP_ALIGN.CENTER

gestures = [
    ("👋 Open Palm", "Start Slideshow"),
    ("✊ Fist", "Stop / Exit"),
    ("☝️ Point (Index)", "Next Slide"),
    ("✌️ Victory (Index+Middle)", "Previous Slide"),
    ("👍 Thumbs Up", "Next Slide"),
    ("👌 OK (Thumb+Index)", "Previous Slide"),
    ("🤙 Call Me (Thumb+Pinky)", "Toggle Pointer"),
    ("🤟 Three Fingers", "Black Screen"),
    ("🖖 Four Fingers", "White Screen"),
]

y_offset = 1.5
for gesture, action in gestures:
    txBox_g = slide2.shapes.add_textbox(Inches(1.5), Inches(y_offset), Inches(10), Inches(0.6))
    tf_g = txBox_g.text_frame
    p_g = tf_g.paragraphs[0]
    p_g.text = f"{gesture}  →  {action}"
    p_g.font.size = Pt(20)
    p_g.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    y_offset += 0.6

# Slide 3 - Getting Started
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
bg3 = slide3.background
fill3 = bg3.fill
fill3.solid()
fill3.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

txBox4 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf4 = txBox4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "Getting Started"
p4.font.size = Pt(36)
p4.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
p4.font.bold = True
p4.alignment = PP_ALIGN.CENTER

steps = [
    "1. Ensure your webcam is connected",
    "2. Run: python main.py",
    "3. Show 'Open Palm' to start the slideshow",
    "4. Use gestures to navigate slides",
    "5. Show 'Fist' or press 'q' to exit",
]

y_offset = 1.8
for step in steps:
    txBox_s = slide3.shapes.add_textbox(Inches(2), Inches(y_offset), Inches(9), Inches(0.6))
    tf_s = txBox_s.text_frame
    p_s = tf_s.paragraphs[0]
    p_s.text = step
    p_s.font.size = Pt(22)
    p_s.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    y_offset += 0.7

# Slide 4 - Thank You
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
bg4 = slide4.background
fill4 = bg4.fill
fill4.solid()
fill4.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

txBox5 = slide4.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
tf5 = txBox5.text_frame
tf5.word_wrap = True
p5 = tf5.paragraphs[0]
p5.text = "Thank You!"
p5.font.size = Pt(54)
p5.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
p5.font.bold = True
p5.alignment = PP_ALIGN.CENTER

txBox6 = slide4.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
tf6 = txBox6.text_frame
p6 = tf6.paragraphs[0]
p6.text = "Built with OpenCV, MediaPipe, and PyWin32"
p6.font.size = Pt(20)
p6.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p6.alignment = PP_ALIGN.CENTER

# Save
prs.save("Presentation.pptx")
print("[SUCCESS] Sample presentation created: Presentation.pptx")
print(f"       Slides: {len(prs.slides)}")
print("       Run 'python main.py' to start controlling it with gestures!")

